#!/usr/bin/env python3
"""
Convert HTML Slides to PowerPoint Presentations
Extracts content from HTML slides and creates .pptx files with custom VUT styling
"""

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
import os
from pathlib import Path


from pathlib import Path


# VUT Official Brand Color Palette (from VUT slide template)
VUT_NAVY = RGBColor(30, 58, 95)              # #1e3a5f - Main background navy
VUT_NAVY_DARK = RGBColor(20, 42, 70)         # #142a46 - Darker navy for contrast
VUT_GOLD = RGBColor(201, 152, 74)            # #c9984a - Official VUT gold
VUT_CHARCOAL = RGBColor(61, 61, 61)          # #3d3d3d - Dark gray accents
WHITE = RGBColor(255, 255, 255)              # #ffffff - Text and logo
VUT_LIGHT_GRAY = RGBColor(200, 200, 200)     # #c8c8c8 - Light gray for secondary text


def set_slide_background(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_logo(slide, logo_path):
    """Add VUT logo to top left of slide"""
    try:
        if logo_path and logo_path.exists():
            left = Inches(0.3)
            top = Inches(0.3)
            height = Inches(0.6)
            slide.shapes.add_picture(str(logo_path), left, top, height=height)
    except Exception as e:
        pass  # Skip logo if there's an error


def add_title_slide(prs, title, week_num, logo_path=None):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_background(slide, VUT_NAVY)
    
    # Add logo
    add_logo(slide, logo_path)
    
    # Week number
    week_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(0.8)
    )
    tf = week_box.text_frame
    tf.text = f"Week {week_num}"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.color.rgb = VUT_GOLD
    p.font.bold = True
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3), Inches(9), Inches(2)
    )
    tf = title_box.text_frame
    tf.text = title
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    
    # Course code
    course_box = slide.shapes.add_textbox(
        Inches(1), Inches(5.5), Inches(8), Inches(0.6)
    )
    tf = course_box.text_frame
    tf.text = "OPS3 - Virtualization and Cloud Infrastructure"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = VUT_LIGHT_GRAY


def add_objectives_slide(prs, objectives, logo_path=None):
    """Add a 'What You Will Learn This Week' slide with learning objectives"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_background(slide, VUT_NAVY_DARK)
    
    # Add logo
    add_logo(slide, logo_path)
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.0), Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    tf.text = "What You Will Learn This Week"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.color.rgb = VUT_GOLD
    p.font.bold = True
    
    # Objectives area
    content_box = slide.shapes.add_textbox(
        Inches(1.2), Inches(2.2), Inches(7.6), Inches(4.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    
    # Add objectives as bullets
    for idx, objective in enumerate(objectives[:6]):  # Limit to 6 objectives
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = objective
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.space_before = Pt(10)
        p.space_after = Pt(10)


def add_content_slide(prs, title, content_items, logo_path=None):
    """Add a content slide with title and bullet points, extras go to notes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_background(slide, VUT_NAVY_DARK)
    
    # Add logo
    add_logo(slide, logo_path)
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.0), Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.color.rgb = VUT_GOLD
    p.font.bold = True
    
    # Limit slide content to 5 key points
    slide_items = content_items[:5]
    notes_items = content_items[5:] if len(content_items) > 5 else []
    
    # Content area
    content_top = 1.9
    content_height = 4.8
    
    content_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(content_top), Inches(8.4), Inches(content_height)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    
    # Add slide items
    for idx, item in enumerate(slide_items):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    # Add detailed content to notes
    if notes_items or len(content_items) > 0:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        
        # Add summary
        text_frame.text = f"{title}\n\n"
        
        # Add all items to notes
        text_frame.text += "Key Points:\n"
        for idx, item in enumerate(content_items, 1):
            text_frame.text += f"{idx}. {item}\n"


def add_section_slide(prs, section_title, logo_path=None):
    """Add a section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_background(slide, VUT_NAVY)
    
    # Add logo
    add_logo(slide, logo_path)
    
    # Section title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(2)
    )
    tf = title_box.text_frame
    tf.text = section_title
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(48)
    p.font.color.rgb = VUT_GOLD
    p.font.bold = True


def add_code_slide(prs, title, code_text, logo_path=None):
    """Add a slide with code content"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_background(slide, VUT_NAVY_DARK)
    
    # Add logo
    add_logo(slide, logo_path)
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.0), Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = VUT_GOLD
    p.font.bold = True
    
    # Code box
    code_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.9), Inches(8.6), Inches(4.8)
    )
    tf = code_box.text_frame
    tf.text = code_text[:600]  # Limit code length for readability
    tf.word_wrap = False
    
    for paragraph in tf.paragraphs:
        paragraph.font.name = 'Courier New'
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = VUT_LIGHT_GRAY
    
    # Add full code to notes if truncated
    if len(code_text) > 600:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = f"{title}\n\nFull Code:\n{code_text}"


def parse_html_slides(html_path):
    """Parse HTML slides and extract content"""
    with open(html_path, 'r', encoding='utf-8') as f:
        html_content = f.read()
    
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Extract week number and title from first slide
    title_slide = soup.find('div', class_='title-slide')
    if title_slide:
        week_text = title_slide.find('div', class_='week-number')
        week_num = week_text.get_text().strip().replace('Week ', '') if week_text else '1'
        title = title_slide.find('h1').get_text().strip() if title_slide.find('h1') else 'Presentation'
    else:
        week_num = '1'
        title = 'Presentation'
    
    # Try to find and extract learning objectives from student notes
    objectives = []
    # Look for student notes file in the same directory
    student_notes_path = html_path.parent / f"Week_{week_num}_Student_Notes.html"
    if student_notes_path.exists():
        with open(student_notes_path, 'r', encoding='utf-8') as f:
            notes_content = f.read()
        notes_soup = BeautifulSoup(notes_content, 'html.parser')
        
        # Look for learning objectives section
        for header in notes_soup.find_all(['h2', 'h3']):
            text = header.get_text().lower()
            if any(keyword in text for keyword in ['learning objective', 'objectives', 'learning outcome', 'what you will learn', 'goals']):
                # Get the next sibling content
                next_elem = header.find_next_sibling()
                if next_elem and next_elem.name == 'ul':
                    for li in next_elem.find_all('li', recursive=False):
                        obj_text = li.get_text().strip()
                        if obj_text and len(obj_text) < 300:
                            objectives.append(obj_text)
                elif next_elem and next_elem.name == 'ol':
                    for li in next_elem.find_all('li', recursive=False):
                        obj_text = li.get_text().strip()
                        if obj_text and len(obj_text) < 300:
                            objectives.append(obj_text)
                break
    
    # If no objectives found, create generic ones from main headings
    if not objectives:
        all_slides = soup.find_all('div', class_='slide')
        h2_headers = []
        for slide_div in all_slides:
            if 'title-slide' not in slide_div.get('class', []):
                h2 = slide_div.find('h2')
                if h2:
                    header_text = h2.get_text().strip()
                    if header_text and header_text.lower() not in ['summary', 'conclusion', 'review']:
                        h2_headers.append(header_text)
        
        # Create objectives from unique headers
        seen = set()
        for header in h2_headers[:6]:
            if header not in seen:
                objectives.append(f"Understand {header.lower()}")
                seen.add(header)
    
    # Extract all slides
    slides_data = []
    all_slides = soup.find_all('div', class_='slide')
    
    for slide_div in all_slides:
        slide_info = {
            'type': 'content',
            'title': '',
            'content': []
        }
        
        # Check if it's a title slide
        if 'title-slide' in slide_div.get('class', []):
            continue  # Skip title slides, we'll create our own
        
        # Get the main heading (h2 or h3)
        h2 = slide_div.find('h2')
        h3 = slide_div.find('h3')
        
        if h2:
            slide_info['title'] = h2.get_text().strip()
            # Check if this is a section divider (slide with only h2, no other content)
            if not h3 and not slide_div.find('ul') and not slide_div.find('p') and not slide_div.find('pre'):
                slide_info['type'] = 'section'
        elif h3:
            slide_info['title'] = h3.get_text().strip()
        
        # Extract content
        # Get bullet lists
        ul_elements = slide_div.find_all('ul', recursive=False)
        for ul in ul_elements:
            for li in ul.find_all('li', recursive=False):
                text = li.get_text().strip()
                if text and len(text) < 500:  # Limit text length
                    slide_info['content'].append(text)
        
        # Get paragraphs (if no lists)
        if not slide_info['content']:
            p_elements = slide_div.find_all('p', recursive=False)
            for p in p_elements[:5]:  # Limit to 5 paragraphs
                text = p.get_text().strip()
                if text and len(text) < 500:
                    slide_info['content'].append(text)
        
        # Get code blocks
        pre_elements = slide_div.find_all('pre')
        if pre_elements:
            for pre in pre_elements[:1]:  # Only first code block
                code_text = pre.get_text().strip()
                if code_text:
                    slide_info['type'] = 'code'
                    slide_info['code'] = code_text[:1000]  # Limit code length
        
        # Only add slides with content
        if slide_info['title'] or slide_info['content']:
            slides_data.append(slide_info)
    
    return week_num, title, objectives, slides_data


def create_powerpoint(html_path, output_path, logo_path=None):
    """Create PowerPoint presentation from HTML slides"""
    week_num, title, objectives, slides_data = parse_html_slides(html_path)
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Add title slide
    add_title_slide(prs, title, week_num, logo_path)
    
    # Add learning objectives slide
    if objectives:
        add_objectives_slide(prs, objectives, logo_path)
    
    # Content slides
    for section in slides_data:
        # Section title slide
        if section['type'] == 'section':
            add_section_slide(prs, section['title'], logo_path)
        elif section['type'] == 'code':
            add_code_slide(prs, section['title'], section.get('code', ''), logo_path)
        else:
            # Regular content slide - limit to 5 items per slide
            if section['content']:
                # Split into slides with max 5 items each
                max_items = 5
                content_chunks = [section['content'][i:i+max_items] 
                                for i in range(0, len(section['content']), max_items)]
                
                for idx, chunk in enumerate(content_chunks):
                    slide_title = section['title']
                    if len(content_chunks) > 1:
                        slide_title += f" (Part {idx+1})"
                    add_content_slide(prs, slide_title, chunk, logo_path)
            else:
                # Slide with just title
                add_section_slide(prs, section['title'], logo_path)
    
    # Add summary slide
    add_section_slide(prs, "Summary", logo_path)
    
    # Save presentation
    prs.save(output_path)
    return len(prs.slides)


def process_week(week_dir, week_num, logo_path):
    """Process a single week's HTML slides and create PowerPoint"""
    html_path = week_dir / f"Week_{week_num}_Slides.html"
    pptx_path = week_dir / f"Week_{week_num}_Slides.pptx"
    
    if not html_path.exists():
        print(f"⚠️  Week {week_num}: HTML slides not found at {html_path}")
        return False
    
    try:
        num_slides = create_powerpoint(html_path, pptx_path, logo_path)
        print(f"✅ Week {week_num}: Created PowerPoint with {num_slides} slides -> {pptx_path.name}")
        return True
    except Exception as e:
        print(f"❌ Week {week_num}: Error creating PowerPoint - {str(e)}")
        return False


def main():
    """Main function to process all weeks"""
    base_dir = Path(__file__).parent.parent
    logo_path = base_dir / "ops3_logo.png"
    
    print("=" * 70)
    print("Converting HTML Slides to PowerPoint Presentations")
    print("=" * 70)
    print()
    
    if logo_path.exists():
        print(f"✓ Using logo: {logo_path.name}")
    else:
        print(f"⚠ Logo not found at {logo_path}")
    print()
    
    weeks = [
        ("Week 1 - Introduction to Virtualization", 1),
        ("Week 2 - Virtual Machines", 2),
        ("Week 3 - Virtual Networking and Linux Networking Fundamentals", 3),
        ("Week 4 - Storage and Backup", 4),
        ("Week 5 - Containers and Resource Management", 5),
        ("Week 6 - Proxmox Cluster and High Availability", 6),
        ("Week 7 - Transition to Cloud Computing Concepts", 7),
        ("Week 8 - Cloud Foundation", 8),
        ("Week 9 - Compute Operations", 9),
        ("Week 10 - Storage and Persistence", 10),
        ("Week 11 - Automation and Cloud API", 11),
        ("Week 12 - Final Project and Review", 12),
    ]
    
    success_count = 0
    for week_name, week_num in weeks:
        week_dir = base_dir / week_name
        if week_dir.exists():
            if process_week(week_dir, week_num, logo_path):
                success_count += 1
        else:
            print(f"⚠️  Week {week_num}: Directory not found: {week_name}")
    
    print()
    print("=" * 70)
    print(f"✅ Successfully created {success_count}/{len(weeks)} PowerPoint presentations")
    print("=" * 70)


if __name__ == "__main__":
    main()
